#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
build_pptx.py — 生成《社科实证论文全流程 Workflow》30 页 PPTX

参考 skills/67-econfin-workflow-toolkit 的技能编排，把"从选题到投稿"的
实证研究工作流做成一份可演示、可编辑的 .pptx。

用法:
    python3 build_pptx.py
输出:
    社科实证论文工作流.pptx   (16:9, 30 页)

依赖: python-pptx  (pip install python-pptx)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ----------------------------------------------------------------------------
# 主题色 (学术深蓝 + 暖橙强调 + 中性灰)
# ----------------------------------------------------------------------------
NAVY      = RGBColor(0x1F, 0x3A, 0x5F)   # 主色 深蓝
NAVY_DK   = RGBColor(0x14, 0x28, 0x44)   # 更深
TEAL      = RGBColor(0x2E, 0x86, 0x86)   # 辅助 青
ORANGE    = RGBColor(0xE0, 0x7B, 0x39)   # 强调 暖橙
GOLD      = RGBColor(0xD9, 0xA5, 0x21)   # 强调 金
INK       = RGBColor(0x22, 0x26, 0x33)   # 正文近黑
GREY      = RGBColor(0x5B, 0x63, 0x72)   # 次级灰
LGREY     = RGBColor(0xEC, 0xEF, 0xF3)   # 浅灰底
CARD      = RGBColor(0xF5, 0xF7, 0xFA)   # 卡片底
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LINE      = RGBColor(0xD3, 0xDA, 0xE4)

CJK_FONT  = "Microsoft YaHei"   # Windows 通用; Mac 回退到系统中文字体
EN_FONT   = "Calibri"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

# ----------------------------------------------------------------------------
# 底层辅助函数
# ----------------------------------------------------------------------------
def _set_cjk(run, font_name=CJK_FONT):
    """让中文字符也使用指定字体 (设置 latin / ea / cs typeface)。"""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font_name)


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color, line_color=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line_color is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_color
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    return tb, tf


def para(tf, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         font=CJK_FONT, space_after=6, space_before=0, level=0, italic=False,
         line_spacing=1.1, first=False):
    p = tf.paragraphs[0] if (first and tf.paragraphs[0].text == "" and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.level = level
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    _set_cjk(r, font)
    return p, r


def runs_para(tf, segments, size=18, align=PP_ALIGN.LEFT, space_after=6,
              line_spacing=1.15, first=False, level=0):
    """一段里混排多种样式: segments = [(text, color, bold), ...]"""
    p = tf.paragraphs[0] if (first and tf.paragraphs[0].text == "" and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    p.level = level
    for seg in segments:
        text, color, bold = seg[0], seg[1], seg[2]
        fsize = seg[3] if len(seg) > 3 else size
        r = p.add_run()
        r.text = text
        r.font.size = Pt(fsize)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = CJK_FONT
        _set_cjk(r, CJK_FONT)
    return p


# ----------------------------------------------------------------------------
# 通用页眉 (内容页)
# ----------------------------------------------------------------------------
def content_header(slide, kicker, title, page_no, total=30):
    bg(slide, WHITE)
    # 左侧竖条
    rect(slide, Inches(0), Inches(0), Inches(0.22), EMU_H, NAVY)
    rect(slide, Inches(0.22), Inches(0), Inches(0.06), EMU_H, ORANGE)
    # kicker
    tb, tf = textbox(slide, Inches(0.7), Inches(0.42), Inches(11.6), Inches(0.42))
    para(tf, kicker, size=14, color=ORANGE, bold=True, first=True, space_after=0)
    # 标题
    tb, tf = textbox(slide, Inches(0.7), Inches(0.74), Inches(11.6), Inches(0.78))
    para(tf, title, size=29, color=NAVY, bold=True, first=True, space_after=0)
    # 标题下分隔线
    rect(slide, Inches(0.72), Inches(1.58), Inches(11.9), Pt(2.2), LINE)
    # 页码
    tb, tf = textbox(slide, Inches(12.1), Inches(6.95), Inches(1.1), Inches(0.4))
    para(tf, f"{page_no:02d} / {total}", size=11, color=GREY, align=PP_ALIGN.RIGHT,
         first=True, space_after=0)
    # 页脚品牌
    tb, tf = textbox(slide, Inches(0.7), Inches(6.95), Inches(7.0), Inches(0.4))
    para(tf, "社科实证论文 Workflow  ·  econfin-workflow-toolkit", size=10,
         color=GREY, first=True, space_after=0)


# ----------------------------------------------------------------------------
# 卡片: 标题 + 描述 + (可选)标签
# ----------------------------------------------------------------------------
def card(slide, l, t, w, h, title, body_lines, accent=NAVY, tag=None,
         title_size=16, body_size=12.5, fill=CARD):
    rect(slide, l, t, w, h, fill, line_color=LINE, line_w=Pt(1))
    rect(slide, l, t, Inches(0.09), h, accent)
    pad = Inches(0.22)
    tb, tf = textbox(slide, l + pad, t + Inches(0.13), w - pad - Inches(0.15),
                     h - Inches(0.2))
    if tag:
        para(tf, tag, size=10.5, color=accent, bold=True, first=True, space_after=2)
        para(tf, title, size=title_size, color=NAVY, bold=True, space_after=5)
    else:
        para(tf, title, size=title_size, color=NAVY, bold=True, first=True, space_after=5)
    for ln in body_lines:
        para(tf, ln, size=body_size, color=GREY, space_after=3, line_spacing=1.12)
    return tf


def chip(slide, l, t, w, h, text, fill, txt_color=WHITE, size=12, bold=True):
    sp = rect(slide, l, t, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try:
        sp.adjustments[0] = 0.28
    except Exception:
        pass
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = txt_color
    r.font.name = CJK_FONT; _set_cjk(r)
    return sp


def arrow(slide, l, t, w, h, color=ORANGE):
    return rect(slide, l, t, w, h, color, shape=MSO_SHAPE.CHEVRON)


# ============================================================================
# SLIDE 1 — 封面
# ============================================================================
def slide_cover():
    s = add_slide()
    bg(s, NAVY)
    # 背景装饰块
    rect(s, Inches(9.6), Inches(0), Inches(3.733), EMU_H, NAVY_DK)
    rect(s, Inches(9.6), Inches(0), Inches(0.06), EMU_H, ORANGE)
    # 右侧流程徽标 (竖排阶段点)
    stages = ["选题", "设计", "数据", "估计", "呈现", "写作", "评审", "投稿"]
    y = Inches(0.95)
    for i, st in enumerate(stages):
        cx = Inches(10.25)
        dot = rect(s, cx, y, Inches(0.34), Inches(0.34), ORANGE if i in (3,) else TEAL,
                   shape=MSO_SHAPE.OVAL)
        tb, tf = textbox(s, cx + Inches(0.5), y - Inches(0.02), Inches(2.4), Inches(0.4),
                         anchor=MSO_ANCHOR.MIDDLE)
        para(tf, st, size=15, color=WHITE, bold=True, first=True, space_after=0)
        if i < len(stages) - 1:
            rect(s, cx + Inches(0.15), y + Inches(0.34), Pt(2), Inches(0.39), TEAL)
        y += Inches(0.73)

    # 主标题区
    tb, tf = textbox(s, Inches(0.9), Inches(1.9), Inches(8.4), Inches(2.6))
    para(tf, "社科实证论文", size=54, color=WHITE, bold=True, first=True, space_after=2)
    para(tf, "全流程 Workflow", size=54, color=ORANGE, bold=True, space_after=14)
    para(tf, "从一个想法到一篇可投稿论文 · 用 AI 技能编排的端到端实证研究流水线",
         size=17, color=RGBColor(0xC6, 0xD2, 0xE2), space_after=0, line_spacing=1.25)

    rect(s, Inches(0.95), Inches(5.0), Inches(2.6), Pt(3), ORANGE)
    tb, tf = textbox(s, Inches(0.9), Inches(5.25), Inches(8.4), Inches(1.6))
    para(tf, "演示方法：双重差分 DiD（配套可运行 Notebook）", size=15,
         color=WHITE, bold=True, first=True, space_after=6)
    para(tf, "基于 econfin-workflow-toolkit · 47 个技能 · 8 个阶段",
         size=13, color=RGBColor(0x9F, 0xB0, 0xC8), space_after=0)


# ============================================================================
# SLIDE 2 — 目录
# ============================================================================
def slide_agenda():
    s = add_slide()
    content_header(s, "AGENDA", "目录 · 七个板块 + 一次实操演示", 2)
    items = [
        ("01", "为什么需要『工作流』思维", "实证研究的痛点与流水线总览", NAVY),
        ("02", "选题与研究设计", "idea-finder · novelty-check · proposal", TEAL),
        ("03", "数据获取与清洗", "data-fetcher · data-cleaning", TEAL),
        ("04", "计量识别与估计", "OLS / Panel / IV / DiD / RDD / SCM", ORANGE),
        ("05", "聚焦 DiD（本次演示）", "平行趋势 · 事件研究 · 交错 DiD · 稳健性", ORANGE),
        ("06", "结果呈现", "table · figure 出版级表图", TEAL),
        ("07", "写作 · 评审 · 投稿", "paper-pipeline · referee · submission", NAVY),
    ]
    col_w = Inches(5.85)
    x0 = Inches(0.75); x1 = Inches(6.85)
    y = Inches(1.95)
    for i, (num, t, sub, ac) in enumerate(items):
        col = i % 2
        row = i // 2
        x = x0 if col == 0 else x1
        yy = y + Inches(row * 1.18)
        rect(s, x, yy, col_w, Inches(1.0), CARD, line_color=LINE)
        rect(s, x, yy, Inches(0.09), Inches(1.0), ac)
        tb, tf = textbox(s, x + Inches(0.22), yy + Inches(0.08), Inches(1.0), Inches(0.85),
                         anchor=MSO_ANCHOR.MIDDLE)
        para(tf, num, size=30, color=ac, bold=True, first=True, space_after=0)
        tb, tf = textbox(s, x + Inches(1.25), yy + Inches(0.12), col_w - Inches(1.4),
                         Inches(0.8), anchor=MSO_ANCHOR.MIDDLE)
        para(tf, t, size=16.5, color=NAVY, bold=True, first=True, space_after=2)
        para(tf, sub, size=11.5, color=GREY, space_after=0)
    # 第 8 格: 演示
    x, yy = x1, y + Inches(3 * 1.18)
    rect(s, x, yy, col_w, Inches(1.0), NAVY)
    rect(s, x, yy, Inches(0.09), Inches(1.0), ORANGE)
    tb, tf = textbox(s, x + Inches(0.22), yy + Inches(0.08), Inches(1.0), Inches(0.85),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "▶", size=26, color=ORANGE, bold=True, first=True, space_after=0)
    tb, tf = textbox(s, x + Inches(1.25), yy + Inches(0.12), col_w - Inches(1.4),
                     Inches(0.8), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "实操：did_demo.ipynb", size=16.5, color=WHITE, bold=True, first=True, space_after=2)
    para(tf, "一键跑通一个完整 DiD 分析", size=11.5, color=RGBColor(0xB9,0xC6,0xDA), space_after=0)


# ============================================================================
# 分节页
# ============================================================================
def slide_section(num, title_cn, title_en, sub, page):
    s = add_slide()
    bg(s, NAVY)
    rect(s, Inches(0), Inches(0), EMU_W, EMU_H, NAVY)
    rect(s, Inches(0), Inches(6.9), EMU_W, Inches(0.6), NAVY_DK)
    # 大号编号
    tb, tf = textbox(s, Inches(0.85), Inches(1.55), Inches(5.5), Inches(4))
    para(tf, num, size=200, color=NAVY_DK, bold=True, first=True, space_after=0)
    # 文案
    rect(s, Inches(4.8), Inches(2.55), Inches(0.12), Inches(2.2), ORANGE)
    tb, tf = textbox(s, Inches(5.15), Inches(2.55), Inches(7.4), Inches(2.6))
    para(tf, title_cn, size=46, color=WHITE, bold=True, first=True, space_after=4)
    para(tf, title_en, size=18, color=ORANGE, bold=True, space_after=16)
    para(tf, sub, size=15, color=RGBColor(0xC6,0xD2,0xE2), space_after=0, line_spacing=1.3)
    tb, tf = textbox(s, Inches(12.1), Inches(6.95), Inches(1.1), Inches(0.4))
    para(tf, f"{page:02d} / 30", size=11, color=RGBColor(0x7A,0x8B,0xA6),
         align=PP_ALIGN.RIGHT, first=True, space_after=0)


# ============================================================================
# SLIDE 4 — 实证研究的痛点
# ============================================================================
def slide_pain():
    s = add_slide()
    content_header(s, "1 · 为什么需要工作流", "实证研究真正难的，不是某一步，而是『串起来』", 4)
    pains = [
        ("选题悬浮", "想法多但说不清新意与贡献，不知道能投哪", ORANGE),
        ("数据黑箱", "清洗口径无记录，半年后自己都复现不了", TEAL),
        ("方法误用", "交错 DiD 直接上 TWFE，识别假设不自检", ORANGE),
        ("表图粗糙", "回归表手工拼贴，格式与期刊要求不符", TEAL),
        ("写作返工", "结构反复推翻，润色、自审、改格式来回拉扯", ORANGE),
        ("投稿低效", "期刊匹配靠猜，cover letter 与格式临时赶", TEAL),
    ]
    x0 = Inches(0.75); y0 = Inches(1.95)
    cw, ch = Inches(3.95), Inches(1.45)
    gap = Inches(0.18)
    for i, (t, d, ac) in enumerate(pains):
        col = i % 3; row = i // 3
        x = x0 + col * (cw + gap)
        y = y0 + row * (ch + gap)
        rect(s, x, y, cw, ch, CARD, line_color=LINE)
        rect(s, x, y, cw, Inches(0.07), ac)
        tb, tf = textbox(s, x + Inches(0.22), y + Inches(0.16), cw - Inches(0.4), ch - Inches(0.3))
        para(tf, t, size=16.5, color=NAVY, bold=True, first=True, space_after=4)
        para(tf, d, size=12.5, color=GREY, space_after=0, line_spacing=1.18)
    # 底部结论条
    y = Inches(5.35)
    rect(s, x0, y, Inches(11.85), Inches(1.25), NAVY)
    rect(s, x0, y, Inches(0.12), Inches(1.25), ORANGE)
    tb, tf = textbox(s, x0 + Inches(0.35), y + Inches(0.16), Inches(11.3), Inches(1.0),
                     anchor=MSO_ANCHOR.MIDDLE)
    runs_para(tf, [("破局点：", ORANGE, True, 17),
                   ("把研究拆成可复现、可交接、可自检的 ", WHITE, False, 17),
                   ("8 个阶段", ORANGE, True, 17),
                   ("，每一步配一个专职技能 —— 流程显性化，质量才可控。", WHITE, False, 17)],
              first=True, space_after=4)
    para(tf, "econfin-workflow-toolkit 的设计哲学：高层 orchestrator（如 paper-pipeline）按固定顺序调度单点技能。",
         size=12, color=RGBColor(0xB9,0xC6,0xDA), space_after=0)


# ============================================================================
# SLIDE 5 — 八阶段工作流总览（流程图）
# ============================================================================
def slide_pipeline_overview():
    s = add_slide()
    content_header(s, "1 · 工作流总览", "八个阶段，一条主线：想法 → 可投稿论文", 5)
    stages = [
        ("①", "选题", "Ideation", TEAL),
        ("②", "设计", "Design", TEAL),
        ("③", "数据", "Data", GOLD),
        ("④", "估计", "Estimation", ORANGE),
        ("⑤", "呈现", "Tables/Figures", GOLD),
        ("⑥", "写作", "Writing", TEAL),
        ("⑦", "评审", "Review", TEAL),
        ("⑧", "投稿", "Submission", NAVY),
    ]
    # 上排流程带
    x = Inches(0.7); y = Inches(2.05)
    bw, bh = Inches(1.36), Inches(1.5)
    gap = Inches(0.13)
    for i, (n, cn, en, ac) in enumerate(stages):
        rect(s, x, y, bw, bh, CARD, line_color=LINE)
        rect(s, x, y, bw, Inches(0.42), ac)
        tb, tf = textbox(s, x, y + Inches(0.02), bw, Inches(0.4), anchor=MSO_ANCHOR.MIDDLE)
        para(tf, f"{n} {cn}", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
             first=True, space_after=0)
        tb, tf = textbox(s, x, y + Inches(0.5), bw, Inches(0.95), anchor=MSO_ANCHOR.TOP)
        para(tf, en, size=10.5, color=GREY, align=PP_ALIGN.CENTER, first=True, space_after=0)
        if i < len(stages) - 1:
            arrow(s, x + bw - Inches(0.02), y + Inches(0.52), Inches(0.17), Inches(0.45), LINE)
        x += bw + gap
    # 中部说明
    tb, tf = textbox(s, Inches(0.75), Inches(3.85), Inches(11.85), Inches(0.5))
    runs_para(tf, [("贯穿全程：", NAVY, True, 14),
                   ("web-research / arxiv 查文献 · stata / stats 计量底座 · reference-verify 引用核查 · markitdown 文档互转",
                    GREY, False, 13.5)], first=True, space_after=0)
    # 下排: 关键产出物
    outs = [
        ("研究问题 + 贡献点", TEAL),
        ("识别策略 + 期刊画像", TEAL),
        ("可复现分析数据集", GOLD),
        ("ATT 估计 + 稳健性", ORANGE),
        ("出版级表 & 图", GOLD),
        ("结构完整初稿", TEAL),
        ("审稿意见 + 回应", TEAL),
        ("投稿包 + cover letter", NAVY),
    ]
    x = Inches(0.7); y = Inches(4.5)
    for i, (o, ac) in enumerate(outs):
        rect(s, x, y, bw, Inches(0.92), WHITE, line_color=ac, line_w=Pt(1.4))
        tb, tf = textbox(s, x + Inches(0.05), y + Inches(0.05), bw - Inches(0.1), Inches(0.82),
                         anchor=MSO_ANCHOR.MIDDLE)
        para(tf, "产出", size=9, color=ac, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=2)
        para(tf, o, size=10, color=INK, align=PP_ALIGN.CENTER, space_after=0, line_spacing=1.05)
        x += bw + gap
    # 底注
    tb, tf = textbox(s, Inches(0.75), Inches(5.75), Inches(11.85), Inches(0.9))
    para(tf, "每个阶段都有明确『输入 → 处理 → 产出』，上一阶段的产出即下一阶段的输入。"
             "这让协作、复现与质量门禁成为可能 —— 这正是『工作流』相对『单点提示词』的根本优势。",
         size=12.5, color=GREY, first=True, space_after=0, line_spacing=1.25)


# ============================================================================
# SLIDE 6 — 工具箱全景
# ============================================================================
def slide_toolkit_map():
    s = add_slide()
    content_header(s, "1 · 工具箱全景", "47 个技能，按阶段编排成一张地图", 6)
    groups = [
        ("选题 & 设计", TEAL, ["econfin-idea-finder", "econfin-proposal", "novelty-check",
                              "significance-search", "journal-digest", "master-thesis-review"]),
        ("数据", GOLD, ["data-fetcher", "data-cleaning"]),
        ("计量估计", ORANGE, ["ols-regression", "panel-data", "iv-estimation", "did-analysis",
                            "rdd-analysis", "synthetic-control", "time-series", "ml-causal",
                            "stata", "stats"]),
        ("表 & 图", GOLD, ["table", "figure"]),
        ("写作 & 润色", TEAL, ["paper-writer", "paper-style", "paper-polish",
                             "paper-self-revise", "paper-pipeline", "readability"]),
        ("评审 & 引用", NAVY, ["referee-report", "paper-referee-revise", "reference-verify"]),
        ("投稿", NAVY, ["paper-submission"]),
        ("写作辅助 / 转换 / 检索", GREY, ["marp-slides-creator", "chinese-ppt", "md-to-docx",
                                    "markitdown", "web-research", "arxiv", "agent-browser",
                                    "fix-chinese"]),
    ]
    x0 = Inches(0.72); y0 = Inches(1.9)
    cw = Inches(3.86); gap_x = Inches(0.12)
    positions = [(0,0),(1,0),(2,0),(0,1),(1,1),(0,2),(1,2),(2,1)]
    # 自定义布局: 3 列
    col_x = [x0, x0 + cw + gap_x, x0 + 2*(cw + gap_x)]
    col_y = [y0, y0, y0]
    layout = {0:0, 4:0, 7:0,   2:1, 3:1, 5:1,  1:2, 6:2}
    heights = {}
    for i, (name, ac, skills) in enumerate(groups):
        col = layout[i]
        x = col_x[col]
        y = col_y[col]
        rows = (len(skills) + 1) // 2
        h = Inches(0.62 + rows * 0.34)
        rect(s, x, y, cw, h, CARD, line_color=LINE)
        rect(s, x, y, cw, Inches(0.36), ac)
        tb, tf = textbox(s, x + Inches(0.15), y + Inches(0.0), cw - Inches(0.3), Inches(0.36),
                         anchor=MSO_ANCHOR.MIDDLE)
        para(tf, f"{name}  ·  {len(skills)}", size=12.5, color=WHITE, bold=True,
             first=True, space_after=0)
        tb, tf = textbox(s, x + Inches(0.14), y + Inches(0.42), cw - Inches(0.26), h - Inches(0.5))
        # 两两一行
        line = ""
        for j, sk in enumerate(skills):
            line += "• " + sk + "   "
            if j % 2 == 1 or j == len(skills) - 1:
                para(tf, line.strip(), size=10.3, color=INK, space_after=2, line_spacing=1.05)
                line = ""
        col_y[col] = y + h + Inches(0.12)
    # 右下角说明卡
    tb, tf = textbox(s, col_x[2], col_y[2] + Inches(0.05), cw, Inches(1.4))
    para(tf, "编排方式", size=12.5, color=ORANGE, bold=True, first=True, space_after=4)
    para(tf, "高层 orchestrator（paper-pipeline）按固定顺序调用单点技能；"
             "你也可以单独调用任意一个技能。", size=10.8, color=GREY, space_after=0, line_spacing=1.2)


# ============================================================================
# SLIDE 8 — 选题三件套
# ============================================================================
def slide_ideation():
    s = add_slide()
    content_header(s, "2 · 选题", "把模糊的兴趣，磨成『有新意、可识别、能投稿』的问题", 8)
    cards = [
        ("econfin-idea-finder", "想法生成", TEAL,
         ["从文献缺口 / 政策事件 / 新数据出发", "批量产生候选研究问题", "标注潜在识别策略与数据来源"]),
        ("novelty-check", "新意核查", ORANGE,
         ["对照已有文献检索相似研究", "定位你的边际贡献在哪", "避免『重复造轮子』"]),
        ("significance-search", "意义检索", GOLD,
         ["评估问题的理论与现实重要性", "对接政策 / 行业关切", "为引言的『so what』备料"]),
    ]
    x = Inches(0.75); y = Inches(2.0); cw = Inches(3.86); gap = Inches(0.13)
    for name, role, ac, lines in cards:
        card(s, x, y, cw, Inches(2.9), name, lines, accent=ac, tag=role,
             title_size=15.5, body_size=12.5)
        x += cw + gap
    # 下方流向
    y2 = Inches(5.2)
    rect(s, Inches(0.75), y2, Inches(11.85), Inches(1.35), CARD, line_color=LINE)
    rect(s, Inches(0.75), y2, Inches(0.12), Inches(1.35), NAVY)
    tb, tf = textbox(s, Inches(1.05), y2 + Inches(0.16), Inches(11.4), Inches(1.05),
                     anchor=MSO_ANCHOR.MIDDLE)
    runs_para(tf, [("输出 → ", NAVY, True, 14),
                   ("一句话研究问题 + 候选识别策略 + 3–5 篇最相关文献 + 边际贡献陈述",
                    INK, False, 14)], first=True, space_after=6)
    para(tf, "这份『选题卡片』直接喂给下一步 econfin-proposal，作为研究设计的输入。",
         size=12.5, color=GREY, space_after=0)


# ============================================================================
# SLIDE 9 — 立项与期刊画像
# ============================================================================
def slide_proposal():
    s = add_slide()
    content_header(s, "2 · 研究设计", "立项 = 把『问题』翻译成『可执行的识别 + 数据 + 期刊』", 9)
    left = [
        ("econfin-proposal · 研究计划书", NAVY,
         ["研究问题与假设 H1/H2", "识别策略：DiD / IV / RDD / SCM 选型",
          "样本、变量与数据来源清单", "实证设计 + 预期结果 + 稳健性预案"]),
        ("journal-digest · 期刊画像", TEAL,
         ["目标期刊近年选题与方法偏好", "字数、结构、表图规范", "为后续 paper-style 提供模板"]),
    ]
    x = Inches(0.75); y = Inches(2.0); cw = Inches(5.85)
    for name, ac, lines in left:
        card(s, x, y, cw, Inches(2.25), name, lines, accent=ac, title_size=15, body_size=12.3)
        y += Inches(2.42)
    # 右侧: 设计自检清单
    rx = Inches(6.85); ry = Inches(2.0)
    rect(s, rx, ry, Inches(5.75), Inches(4.67), NAVY)
    rect(s, rx, ry, Inches(0.12), Inches(4.67), ORANGE)
    tb, tf = textbox(s, rx + Inches(0.35), ry + Inches(0.22), Inches(5.2), Inches(4.3))
    para(tf, "立项阶段必答的 6 个问题", size=16, color=ORANGE, bold=True, first=True, space_after=10)
    qs = ["1. 因果问题是什么？反事实是什么？",
          "2. 变异来自哪里？为什么外生？",
          "3. 用什么识别策略？关键假设可检验吗？",
          "4. 数据能否支撑该设计（频率/层级/样本量）？",
          "5. 主要威胁与稳健性检验清单？",
          "6. 结果若显著 / 不显著，分别意味着什么？"]
    for q in qs:
        para(tf, q, size=13, color=WHITE, space_after=9, line_spacing=1.15)


# ============================================================================
# SLIDE 11 — 数据获取
# ============================================================================
def slide_data_fetch():
    s = add_slide()
    content_header(s, "3 · 数据", "data-fetcher：把分散的原始数据，变成一张主表", 11)
    tb, tf = textbox(s, Inches(0.75), Inches(1.85), Inches(11.8), Inches(0.6))
    para(tf, "面板/横截面研究的第一道工序：定位来源 → 抓取 → 对齐键 → 合并 → 留痕。",
         size=14, color=GREY, first=True, space_after=0)
    sources = [
        ("公开数据库", TEAL, "FRED · World Bank · Wind/CSMAR 导出 · 统计年鉴"),
        ("文件导入", GOLD, "CSV / Excel / dta / parquet 统一读入"),
        ("网络检索", ORANGE, "web-research · arxiv · agent-browser 抓取补充变量"),
        ("API 接口", NAVY, "宏观、金融、专利等结构化接口拉取"),
    ]
    x = Inches(0.75); y = Inches(2.55); cw = Inches(2.9); gap = Inches(0.08)
    for t, ac, d in sources:
        rect(s, x, y, cw, Inches(1.7), CARD, line_color=LINE)
        rect(s, x, y, cw, Inches(0.5), ac)
        tb, tf = textbox(s, x, y + Inches(0.04), cw, Inches(0.45), anchor=MSO_ANCHOR.MIDDLE)
        para(tf, t, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=0)
        tb, tf = textbox(s, x + Inches(0.18), y + Inches(0.62), cw - Inches(0.34), Inches(1.0))
        para(tf, d, size=12, color=GREY, first=True, space_after=0, line_spacing=1.2)
        x += cw + gap
    # 合并键提示
    y2 = Inches(4.55)
    rect(s, Inches(0.75), y2, Inches(11.85), Inches(2.0), NAVY)
    rect(s, Inches(0.75), y2, Inches(0.12), Inches(2.0), ORANGE)
    tb, tf = textbox(s, Inches(1.1), y2 + Inches(0.2), Inches(11.2), Inches(1.7))
    para(tf, "合并的关键：键（key）必须干净且唯一", size=15.5, color=ORANGE, bold=True,
         first=True, space_after=8)
    runs_para(tf, [("• 实体键：", WHITE, True, 12.5), ("公司代码 / 个体 ID / 地区码 —— 跨源统一编码，处理历史变更", RGBColor(0xC6,0xD2,0xE2), False, 12.5)], space_after=5)
    runs_para(tf, [("• 时间键：", WHITE, True, 12.5), ("year / quarter / month —— 频率对齐，注意财政年 vs 自然年", RGBColor(0xC6,0xD2,0xE2), False, 12.5)], space_after=5)
    runs_para(tf, [("• 留痕：", WHITE, True, 12.5), ("每一次抓取与合并记录来源、口径、时间戳 —— 半年后还能复现", RGBColor(0xC6,0xD2,0xE2), False, 12.5)], space_after=0)


# ============================================================================
# SLIDE 12 — 数据清洗清单
# ============================================================================
def slide_data_clean():
    s = add_slide()
    content_header(s, "3 · 数据清洗", "data-cleaning：从『能跑』到『可信』之间的全部脏活", 12)
    items = [
        ("缺失值", TEAL, "标记 vs 插补 vs 删除；区分 MCAR/MAR/MNAR"),
        ("异常值", GOLD, "缩尾 winsorize（1%/99%）或截尾，留对照"),
        ("变量构造", ORANGE, "对数、增长率、比率、标准化、滞后项"),
        ("面板结构", NAVY, "平衡性检查、重复键、进入/退出样本"),
        ("单位与量纲", TEAL, "货币/通胀平减、百万 vs 元、统一计量"),
        ("分类与编码", GOLD, "行业/地区码统一、哑变量、固定效应分组"),
    ]
    x0 = Inches(0.75); y0 = Inches(1.95); cw = Inches(3.86); ch = Inches(1.35)
    gx = Inches(0.13); gy = Inches(0.18)
    for i, (t, ac, d) in enumerate(items):
        col = i % 3; row = i // 3
        x = x0 + col * (cw + gx); y = y0 + row * (ch + gy)
        rect(s, x, y, cw, ch, CARD, line_color=LINE)
        rect(s, x, y, Inches(0.09), ch, ac)
        tb, tf = textbox(s, x + Inches(0.22), y + Inches(0.14), cw - Inches(0.35), ch - Inches(0.25))
        para(tf, t, size=15.5, color=NAVY, bold=True, first=True, space_after=4)
        para(tf, d, size=12, color=GREY, space_after=0, line_spacing=1.18)
    # 底部金句
    y = Inches(5.45)
    rect(s, x0, y, Inches(11.85), Inches(1.1), CARD, line_color=LINE)
    rect(s, x0, y, Inches(0.12), Inches(1.1), ORANGE)
    tb, tf = textbox(s, x0 + Inches(0.35), y + Inches(0.1), Inches(11.3), Inches(0.9),
                     anchor=MSO_ANCHOR.MIDDLE)
    runs_para(tf, [("原则：", ORANGE, True, 15),
                   ("每一步清洗都写成可重跑的脚本 + 一行注释说明『为什么这么处理』。"
                    "清洗日志本身就是论文附录的数据章节。", INK, False, 14)],
              first=True, space_after=0, line_spacing=1.25)


# ============================================================================
# SLIDE 14 — 因果识别策略地图
# ============================================================================
def slide_method_map():
    s = add_slide()
    content_header(s, "4 · 计量识别", "八种识别策略：先有『干净的变异』，再谈回归", 14)
    rows = [
        ("OLS / 多元回归", "ols-regression", "控制混淆后的条件相关", "选择性偏误难排除", GREY),
        ("面板固定效应", "panel-data", "个体/时间不可观测异质性", "时变混淆仍威胁", TEAL),
        ("工具变量 IV", "iv-estimation", "内生性 / 反向因果", "弱工具 & 排他性", GOLD),
        ("双重差分 DiD", "did-analysis", "政策冲击 / 自然实验", "平行趋势 & 交错偏误", ORANGE),
        ("断点回归 RDD", "rdd-analysis", "阈值规则分配处理", "带宽 & 操纵", TEAL),
        ("合成控制 SCM", "synthetic-control", "单一处理单元（地区/国家）", "donor 池 & 拟合", GOLD),
        ("时间序列", "time-series", "宏观/单变量动态", "平稳性 & 结构突变", TEAL),
        ("机器学习因果", "ml-causal", "异质处理效应 / 高维控制", "正则化 & 推断", ORANGE),
    ]
    # 表头
    x0 = Inches(0.75); y0 = Inches(1.9)
    cols = [Inches(2.6), Inches(2.5), Inches(3.45), Inches(3.3)]
    headers = ["方法", "对应技能", "适用场景", "核心威胁"]
    hx = x0
    rect(s, x0, y0, sum(cols, Emu(0)), Inches(0.5), NAVY)
    for c, htxt in zip(cols, headers):
        tb, tf = textbox(s, hx + Inches(0.1), y0 + Inches(0.02), c - Inches(0.15), Inches(0.46),
                         anchor=MSO_ANCHOR.MIDDLE)
        para(tf, htxt, size=13, color=WHITE, bold=True, first=True, space_after=0)
        hx += c
    rh = Inches(0.55)
    for i, (m, sk, use, threat, ac) in enumerate(rows):
        y = y0 + Inches(0.5) + i * rh
        fill = WHITE if i % 2 == 0 else LGREY
        rect(s, x0, y, sum(cols, Emu(0)), rh, fill, line_color=LINE)
        rect(s, x0, y, Inches(0.07), rh, ac)
        cellx = x0
        vals = [m, sk, use, threat]
        sizes = [12.5, 11.5, 11.5, 11.5]
        bolds = [True, False, False, False]
        colors = [NAVY, ac, INK, GREY]
        for c, v, sz, bd, cl in zip(cols, vals, sizes, bolds, colors):
            tb, tf = textbox(s, cellx + Inches(0.14), y, c - Inches(0.2), rh, anchor=MSO_ANCHOR.MIDDLE)
            para(tf, v, size=sz, color=cl, bold=bd, first=True, space_after=0, line_spacing=1.0)
            cellx += c
    # 右注
    tb, tf = textbox(s, Inches(0.75), Inches(6.85), Inches(11.6), Inches(0.4))
    para(tf, "底座技能：stata（执行）· stats（检验与诊断）贯穿所有方法。",
         size=11.5, color=GREY, first=True, space_after=0)


# ============================================================================
# SLIDE 15 — 方法选择决策树
# ============================================================================
def slide_decision_tree():
    s = add_slide()
    content_header(s, "4 · 方法选型", "一棵决策树：处理是怎么分配的？", 15)
    # 根
    chip(s, Inches(5.2), Inches(1.85), Inches(3.0), Inches(0.6),
         "处理（policy/treatment）如何分配？", NAVY, WHITE, size=13)
    branches = [
        ("有明确时点的政策冲击", "→ DiD / 事件研究", ORANGE,
         "多期、有处理组与对照组，且能论证平行趋势"),
        ("由阈值规则决定", "→ RDD 断点回归", TEAL,
         "running variable 跨过 cutoff 即受处理"),
        ("存在外生工具变量", "→ IV / 2SLS", GOLD,
         "工具影响处理、但只通过处理影响结果"),
        ("仅单一处理单元", "→ 合成控制 SCM", TEAL,
         "用对照单元的加权组合构造反事实"),
        ("仅可控混淆 / 观测性", "→ Panel FE / OLS / ML", GREY,
         "靠固定效应与控制变量逼近条件独立"),
    ]
    x = Inches(0.75); y = Inches(2.95)
    cw = Inches(2.32); gap = Inches(0.07)
    for cond, method, ac, note in branches:
        rect(s, x, y, cw, Inches(3.3), CARD, line_color=LINE)
        rect(s, x, y, cw, Inches(0.95), ac)
        tb, tf = textbox(s, x + Inches(0.12), y + Inches(0.08), cw - Inches(0.24), Inches(0.85),
                         anchor=MSO_ANCHOR.MIDDLE)
        para(tf, cond, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True,
             space_after=0, line_spacing=1.1)
        tb, tf = textbox(s, x + Inches(0.14), y + Inches(1.1), cw - Inches(0.28), Inches(2.1))
        para(tf, method, size=14.5, color=ac, bold=True, first=True, space_after=8, line_spacing=1.1)
        para(tf, note, size=11, color=GREY, space_after=0, line_spacing=1.22)
        # 连接线
        rect(s, x + cw/2 - Pt(1), Inches(2.45), Pt(2), Inches(0.5), LINE)
        x += cw + gap
    tb, tf = textbox(s, Inches(0.75), Inches(6.45), Inches(11.8), Inches(0.5))
    runs_para(tf, [("提醒：", ORANGE, True, 12.5),
                   ("方法服从于识别，不是反过来。先问『反事实从哪来』，再选估计量。本次演示走最左侧分支 → DiD。",
                    GREY, False, 12.5)], first=True, space_after=0)


# ============================================================================
# SLIDE 17 — DiD 核心逻辑
# ============================================================================
def slide_did_core():
    s = add_slide()
    content_header(s, "5 · 聚焦 DiD", "双重差分：用『差中之差』剥离掉共同趋势", 17)
    # 左: 2x2 表
    x0 = Inches(0.75); y0 = Inches(2.0)
    rect(s, x0, y0, Inches(5.9), Inches(2.7), CARD, line_color=LINE)
    tb, tf = textbox(s, x0 + Inches(0.25), y0 + Inches(0.15), Inches(5.4), Inches(0.4))
    para(tf, "2×2 设计：四个均值，两次相减", size=15, color=NAVY, bold=True, first=True, space_after=0)
    # mini table
    tx = x0 + Inches(0.32); ty = y0 + Inches(0.75)
    cellw = Inches(1.28); cellh = Inches(0.55)
    grid = [["", "处理前", "处理后", "差(Δ)"],
            ["处理组", "Ȳ T,pre", "Ȳ T,post", "Δ_T"],
            ["对照组", "Ȳ C,pre", "Ȳ C,post", "Δ_C"]]
    for r in range(3):
        for c in range(4):
            cx = tx + c * cellw; cy = ty + r * cellh
            head = (r == 0 or c == 0)
            fill = NAVY if r == 0 else (LGREY if c == 0 else WHITE)
            if r == 0 and c == 0:
                fill = NAVY
            rect(s, cx, cy, cellw, cellh, fill, line_color=LINE)
            tb, tf = textbox(s, cx, cy, cellw, cellh, anchor=MSO_ANCHOR.MIDDLE)
            col = WHITE if r == 0 else (NAVY if c == 0 else INK)
            para(tf, grid[r][c], size=11.5, color=col, bold=head,
                 align=PP_ALIGN.CENTER, first=True, space_after=0)
    tb, tf = textbox(s, x0 + Inches(0.4), y0 + Inches(2.35), Inches(5.2), Inches(0.4))
    runs_para(tf, [("DiD = Δ_T − Δ_C", ORANGE, True, 15),
                   ("  →  剔除共同时间趋势后的处理效应 ATT", GREY, False, 12)],
              first=True, space_after=0)
    # 右: 回归式
    rx = Inches(6.85); ry = Inches(2.0)
    rect(s, rx, ry, Inches(5.75), Inches(2.7), NAVY)
    tb, tf = textbox(s, rx + Inches(0.3), ry + Inches(0.2), Inches(5.2), Inches(2.4))
    para(tf, "回归形式", size=15, color=ORANGE, bold=True, first=True, space_after=8)
    para(tf, "Yᵢₜ = β₀ + β₁·Treatᵢ + β₂·Postₜ", size=15, color=WHITE, bold=True,
         space_after=2, font=EN_FONT)
    para(tf, "          + β₃·(Treatᵢ × Postₜ) + εᵢₜ", size=15, color=WHITE, bold=True,
         space_after=8, font=EN_FONT)
    runs_para(tf, [("β₃ ", ORANGE, True, 15), ("就是 DiD 估计量（ATT）", WHITE, False, 13)],
              space_after=6)
    para(tf, "面板更优：双向固定效应 TWFE，聚类稳健标准误（按处理层级聚类）。",
         size=12, color=RGBColor(0xC6,0xD2,0xE2), space_after=0, line_spacing=1.2)
    # 底: 核心假设
    y2 = Inches(5.0)
    rect(s, x0, y2, Inches(11.85), Inches(1.55), CARD, line_color=LINE)
    rect(s, x0, y2, Inches(0.12), Inches(1.55), ORANGE)
    tb, tf = textbox(s, x0 + Inches(0.35), y2 + Inches(0.15), Inches(11.3), Inches(1.3))
    para(tf, "唯一的核心假设：平行趋势（Parallel Trends）", size=15, color=NAVY, bold=True,
         first=True, space_after=6)
    para(tf, "若没有处理，处理组的结果本应与对照组『平行』地演变。这个反事实不可直接观测 —— "
             "所以我们用处理前的趋势是否平行来『侧面』验证它。", size=13, color=GREY,
         space_after=0, line_spacing=1.25)


# ============================================================================
# SLIDE 18 — 平行趋势与事件研究
# ============================================================================
def slide_event_study():
    s = add_slide()
    content_header(s, "5 · 平行趋势", "事件研究：把单个 Post 拆成『相对时间』动态系数", 18)
    # 左侧文字
    x0 = Inches(0.75); y0 = Inches(1.95)
    rect(s, x0, y0, Inches(5.7), Inches(4.6), CARD, line_color=LINE)
    rect(s, x0, y0, Inches(0.1), Inches(4.6), TEAL)
    tb, tf = textbox(s, x0 + Inches(0.3), y0 + Inches(0.22), Inches(5.25), Inches(4.2))
    para(tf, "做法", size=15, color=TEAL, bold=True, first=True, space_after=5)
    para(tf, "用相对时间哑变量替换单一 Post，以 t = −1 为基准：", size=13, color=INK,
         space_after=6, line_spacing=1.2)
    para(tf, "Yᵢₜ = Σₖ βₖ·1{t−E=k} + αᵢ + δₜ + εᵢₜ", size=13.5, color=NAVY, bold=True,
         space_after=10, font=EN_FONT)
    para(tf, "怎么读这张图", size=15, color=TEAL, bold=True, space_after=5)
    for ln in ["• 处理前系数 ≈ 0 → 支持平行趋势",
               "• 处理前联合 F 检验：所有 pre 系数 = 0",
               "• 处理后系数 → 动态处理效应轨迹",
               "• 警惕预期效应（t−2 起就偏离 0）",
               "• 注意 Ashenfelter's dip（处理前下陷）"]:
        para(tf, ln, size=12.5, color=GREY, space_after=5, line_spacing=1.15)
    # 右侧: 简易事件研究示意图（用形状画）
    gx = Inches(6.75); gy = Inches(1.95); gw = Inches(5.85); gh = Inches(4.6)
    rect(s, gx, gy, gw, gh, WHITE, line_color=LINE)
    tb, tf = textbox(s, gx + Inches(0.2), gy + Inches(0.1), gw - Inches(0.4), Inches(0.4))
    para(tf, "事件研究示意（pre≈0，post 跳升）", size=12.5, color=NAVY, bold=True,
         first=True, space_after=0)
    # 坐标
    base_y = gy + Inches(2.7)   # 0 线
    left_x = gx + Inches(0.7)
    axis_w = gw - Inches(1.2)
    rect(s, left_x, base_y, axis_w, Pt(1.5), GREY)          # x 轴 (0线)
    rect(s, left_x, gy + Inches(0.6), Pt(1.5), Inches(3.6), GREY)  # y 轴
    # 处理时点竖虚线
    evx = left_x + axis_w * 0.5
    rect(s, evx, gy + Inches(0.6), Pt(1.5), Inches(3.6), ORANGE)
    # 点: 相对时间 -3..3
    pts = [(-3,0.02),(-2,-0.03),(-1,0.0),(0,0.10),(1,0.16),(2,0.18),(3,0.19)]
    n = len(pts)
    for k,(rt,val) in enumerate(pts):
        px = left_x + axis_w * (k/(n-1))
        py = base_y - Inches(val*9)   # 缩放
        dotc = TEAL if rt < 0 else ORANGE
        d = rect(s, px - Inches(0.07), py - Inches(0.07), Inches(0.14), Inches(0.14), dotc,
                 shape=MSO_SHAPE.OVAL)
        # 误差棒
        rect(s, px - Pt(0.8), py - Inches(0.18), Pt(1.6), Inches(0.36), dotc)
    tb, tf = textbox(s, evx - Inches(0.6), base_y + Inches(0.15), Inches(1.4), Inches(0.3))
    para(tf, "处理时点 E", size=10.5, color=ORANGE, bold=True, align=PP_ALIGN.CENTER,
         first=True, space_after=0)
    tb, tf = textbox(s, gx + Inches(0.2), gy + gh - Inches(0.5), gw - Inches(0.4), Inches(0.4))
    para(tf, "蓝点=处理前（应贴近 0 线）   橙点=处理后（显著为正）", size=10.5, color=GREY,
         first=True, space_after=0)


# ============================================================================
# SLIDE 19 — 交错 DiD 的陷阱
# ============================================================================
def slide_staggered():
    s = add_slide()
    content_header(s, "5 · 交错 DiD", "当各单位在不同时间被处理：别再无脑 TWFE", 19)
    # 警告条
    rect(s, Inches(0.75), Inches(1.85), Inches(11.85), Inches(1.0), NAVY)
    rect(s, Inches(0.75), Inches(1.85), Inches(0.12), Inches(1.0), ORANGE)
    tb, tf = textbox(s, Inches(1.1), Inches(1.95), Inches(11.3), Inches(0.8), anchor=MSO_ANCHOR.MIDDLE)
    runs_para(tf, [("陷阱：", ORANGE, True, 15),
                   ("处理时点交错 + 效应异质 时，标准 TWFE 会用『已处理单位』当对照，"
                    "产生“负权重”，估计量可能严重有偏（甚至变号）。", WHITE, False, 14)],
              first=True, space_after=0, line_spacing=1.2)
    # 现代估计量卡片
    ests = [
        ("Callaway & Sant'Anna", TEAL, "csdid / did", "按『组×时间』估 ATT(g,t)，再聚合；干净对照组"),
        ("Sun & Abraham", GOLD, "sunab / fixest", "交互加权事件研究，修正污染的相对时间系数"),
        ("Borusyak et al. (BJS)", ORANGE, "did_imputation", "插补法：先用未处理拟合反事实，再求残差"),
        ("de Chaisemartin–D'Haultf.", NAVY, "did_multiplegt", "允许处理开关、连续/多值处理"),
    ]
    x0 = Inches(0.75); y0 = Inches(3.1); cw = Inches(5.85); ch = Inches(1.5)
    gx = Inches(0.13); gy = Inches(0.16)
    for i,(name, ac, tool, desc) in enumerate(ests):
        col = i % 2; row = i // 2
        x = x0 + col*(cw+gx); y = y0 + row*(ch+gy)
        rect(s, x, y, cw, ch, CARD, line_color=LINE)
        rect(s, x, y, Inches(0.1), ch, ac)
        tb, tf = textbox(s, x + Inches(0.25), y + Inches(0.12), cw - Inches(0.4), ch - Inches(0.2))
        runs_para(tf, [(name, NAVY, True, 14.5), ("   ["+tool+"]", ac, True, 11)],
                  first=True, space_after=4)
        para(tf, desc, size=12, color=GREY, space_after=0, line_spacing=1.18)
    # 底: 诊断
    y = Inches(6.45)
    tb, tf = textbox(s, Inches(0.75), y, Inches(11.85), Inches(0.5))
    runs_para(tf, [("先诊断后估计：", ORANGE, True, 12.5),
                   ("用 Goodman-Bacon 分解看『坏对照』权重有多大；权重显著就果断换上面任一现代估计量。",
                    GREY, False, 12.5)], first=True, space_after=0)


# ============================================================================
# SLIDE 20 — DiD 稳健性清单
# ============================================================================
def slide_robustness():
    s = add_slide()
    content_header(s, "5 · 稳健性", "审稿人会逐条核对的 DiD 体检清单", 20)
    checks = [
        ("平行趋势检验", TEAL, "事件研究图 + 处理前联合 F 检验；Roth(2022) 功效分析"),
        ("安慰剂处理时点", GOLD, "把处理提前 1–2 期，伪 DiD 应不显著"),
        ("安慰剂处理组", ORANGE, "仅用对照单元造假处理，系数应为零"),
        ("替换对照组", TEAL, "限定更可比的对照，结论应稳定"),
        ("交错偏误自检", NAVY, "Bacon 分解 + CS / SA / BJS 对照主回归"),
        ("聚类与推断", GOLD, "按处理层级聚类；少簇时 wild cluster bootstrap"),
        ("预期 & 溢出", ORANGE, "检验 anticipation；SUTVA / 邻近溢出缓冲带"),
        ("剂量反应", TEAL, "连续处理强度，验证『越强越大』的单调性"),
    ]
    x0 = Inches(0.75); y0 = Inches(1.95); cw = Inches(5.85); ch = Inches(1.02)
    gx = Inches(0.13); gy = Inches(0.13)
    for i,(t, ac, d) in enumerate(checks):
        col = i % 2; row = i // 2
        x = x0 + col*(cw+gx); y = y0 + row*(ch+gy)
        rect(s, x, y, cw, ch, CARD, line_color=LINE)
        # 序号方块
        rect(s, x, y, Inches(0.55), ch, ac)
        tb, tf = textbox(s, x, y, Inches(0.55), ch, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, f"{i+1:02d}", size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
             first=True, space_after=0)
        tb, tf = textbox(s, x + Inches(0.72), y + Inches(0.1), cw - Inches(0.85), ch - Inches(0.2),
                         anchor=MSO_ANCHOR.MIDDLE)
        para(tf, t, size=14, color=NAVY, bold=True, first=True, space_after=2)
        para(tf, d, size=11, color=GREY, space_after=0, line_spacing=1.1)


# ============================================================================
# SLIDE 21 — Notebook 演示说明
# ============================================================================
def slide_notebook():
    s = add_slide()
    content_header(s, "5 · 实操演示", "did_demo.ipynb：6 个 cell 跑通一个完整 DiD", 21)
    tb, tf = textbox(s, Inches(0.75), Inches(1.85), Inches(11.8), Inches(0.55))
    runs_para(tf, [("配套技能：", NAVY, True, 14),
                   ("did-analysis", ORANGE, True, 14),
                   ("   ·   纯 Python（pandas + statsmodels + matplotlib），打开即跑，无需外部数据。",
                    GREY, False, 13.5)], first=True, space_after=0)
    steps = [
        ("①", "生成模拟面板", "200 个体 × 12 年的交错政策；植入已知真实 ATT，便于对照估计是否还原"),
        ("②", "可视化原始趋势", "处理组 vs 对照组的结果均值时间序列，肉眼看『处理后分叉』"),
        ("③", "2×2 / TWFE 基准", "smf.ols 与双向固定效应回归，聚类稳健标准误，读出 β₃"),
        ("④", "事件研究 + 平行趋势", "相对时间系数图 + 处理前联合检验，给出可放进论文的 Figure"),
        ("⑤", "稳健性：安慰剂", "伪处理时点 / 伪处理组，确认主效应不是噪声"),
        ("⑥", "导出表与图", "回归表（文本/LaTeX）与事件研究图 PNG，直接进入 table / figure 阶段"),
    ]
    x0 = Inches(0.75); y0 = Inches(2.5); cw = Inches(5.85); ch = Inches(1.25)
    gx = Inches(0.13); gy = Inches(0.15)
    for i,(n, t, d) in enumerate(steps):
        col = i % 2; row = i // 2
        x = x0 + col*(cw+gx); y = y0 + row*(ch+gy)
        rect(s, x, y, cw, ch, CARD, line_color=LINE)
        rect(s, x, y, Inches(0.7), ch, NAVY if i%2==0 else TEAL)
        tb, tf = textbox(s, x, y, Inches(0.7), ch, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, n, size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=0)
        tb, tf = textbox(s, x + Inches(0.85), y + Inches(0.12), cw - Inches(1.0), ch - Inches(0.2),
                         anchor=MSO_ANCHOR.MIDDLE)
        para(tf, t, size=14, color=NAVY, bold=True, first=True, space_after=2)
        para(tf, d, size=11, color=GREY, space_after=0, line_spacing=1.12)


# ============================================================================
# SLIDE 23 — 表格规范
# ============================================================================
def slide_table():
    s = add_slide()
    content_header(s, "6 · 结果呈现", "table：把回归结果排成『审稿人一眼能读』的表", 23)
    # 左: 规范清单
    x0 = Inches(0.75); y0 = Inches(2.0)
    rect(s, x0, y0, Inches(5.7), Inches(4.55), CARD, line_color=LINE)
    rect(s, x0, y0, Inches(0.1), Inches(4.55), TEAL)
    tb, tf = textbox(s, x0 + Inches(0.3), y0 + Inches(0.2), Inches(5.2), Inches(4.2))
    para(tf, "出版级回归表要素", size=15.5, color=TEAL, bold=True, first=True, space_after=8)
    for ln in ["• 系数 + 括号内（聚类）标准误",
               "• 显著性星标 *** ** *（脚注口径）",
               "• 控制变量 / 固定效应 用『Yes / No』行",
               "• N、R²（within R²）、聚类数",
               "• 列标题=被解释变量或样本切分",
               "• 三线表（booktabs），无竖线",
               "• 题注自解释：读表不读正文也懂"]:
        para(tf, ln, size=13, color=INK, space_after=7, line_spacing=1.15)
    # 右: 迷你示例表
    rx = Inches(6.75); ry = Inches(2.0)
    rect(s, rx, ry, Inches(5.85), Inches(4.55), WHITE, line_color=LINE)
    tb, tf = textbox(s, rx + Inches(0.25), ry + Inches(0.15), Inches(5.4), Inches(0.4))
    para(tf, "示例：DiD 主回归表", size=13, color=NAVY, bold=True, first=True, space_after=0)
    tab = [
        ["", "(1) OLS", "(2) TWFE"],
        ["Treat×Post", "0.182***", "0.171***"],
        ["", "(0.041)", "(0.038)"],
        ["个体固定效应", "No", "Yes"],
        ["年份固定效应", "No", "Yes"],
        ["N", "2,400", "2,400"],
        ["R²", "0.12", "0.78"],
    ]
    tx = rx + Inches(0.3); ty = ry + Inches(0.65)
    colw = [Inches(2.3), Inches(1.5), Inches(1.5)]
    rh = Inches(0.45)
    for r, rowv in enumerate(tab):
        cx = tx
        top = (r == 0)
        for c, val in enumerate(rowv):
            w = colw[c]
            if top:
                rect(s, cx, ty + r*rh, w, rh, NAVY)
            elif r == 1:
                rect(s, cx, ty + r*rh, w, rh, LGREY)
            tb, tf = textbox(s, cx, ty + r*rh, w, rh, anchor=MSO_ANCHOR.MIDDLE)
            col = WHITE if top else (ORANGE if (r in (1,2) and c>0) else INK)
            bd = top or (r==1 and c>0)
            al = PP_ALIGN.LEFT if c==0 else PP_ALIGN.CENTER
            para(tf, val, size=11.5, color=col, bold=bd, align=al, first=True, space_after=0)
            cx += w
    rect(s, tx, ty, sum(colw, Emu(0)), Pt(2), NAVY)
    rect(s, tx, ty + len(tab)*rh, sum(colw, Emu(0)), Pt(2), NAVY)
    tb, tf = textbox(s, rx + Inches(0.3), ry + Inches(4.0), Inches(5.3), Inches(0.5))
    para(tf, "括号内为公司层面聚类稳健标准误。*** p<0.01。", size=10, color=GREY,
         first=True, space_after=0)


# ============================================================================
# SLIDE 24 — 图形规范
# ============================================================================
def slide_figure():
    s = add_slide()
    content_header(s, "6 · 结果呈现", "figure：一张事件研究图，胜过半页文字", 24)
    cards3 = [
        ("事件研究图", ORANGE, ["相对时间系数 + 95% CI", "处理前贴 0 线 = 平行趋势",
                            "DiD 论文的『门面图』"]),
        ("系数对比图", TEAL, ["多估计量并排：TWFE / CS / SA", "coefplot 风格", "展示结论稳健"]),
        ("分布 / 趋势图", GOLD, ["处理组 vs 对照组均值轨迹", "密度 / 箱线对比", "讲清识别直觉"]),
    ]
    x = Inches(0.75); y = Inches(2.0); cw = Inches(3.86); gap = Inches(0.13)
    for t, ac, lines in cards3:
        card(s, x, y, cw, Inches(2.4), t, lines, accent=ac, title_size=15.5, body_size=12.5)
        x += cw + gap
    # 底部规范
    y2 = Inches(4.7)
    rect(s, Inches(0.75), y2, Inches(11.85), Inches(1.85), NAVY)
    rect(s, Inches(0.75), y2, Inches(0.12), Inches(1.85), ORANGE)
    tb, tf = textbox(s, Inches(1.1), y2 + Inches(0.18), Inches(11.2), Inches(1.55))
    para(tf, "图形的『出版级』标准", size=15, color=ORANGE, bold=True, first=True, space_after=7)
    for seg in [
        [("• 矢量优先：", WHITE, True, 12.5), ("PDF / SVG，放大不糊；字号与正文匹配", RGBColor(0xC6,0xD2,0xE2), False, 12.5)],
        [("• 自解释：", WHITE, True, 12.5), ("坐标轴标签 + 单位 + 图例 + 题注，脱离正文也能读懂", RGBColor(0xC6,0xD2,0xE2), False, 12.5)],
        [("• 克制：", WHITE, True, 12.5), ("去网格噪声、去 3D、配色色盲友好；一图一信息", RGBColor(0xC6,0xD2,0xE2), False, 12.5)],
    ]:
        runs_para(tf, seg, space_after=5)


# ============================================================================
# SLIDE 26 — 写作流水线
# ============================================================================
def slide_writing():
    s = add_slide()
    content_header(s, "7 · 写作", "paper-pipeline：初稿之后的『一条龙打磨』", 26)
    tb, tf = textbox(s, Inches(0.75), Inches(1.85), Inches(11.8), Inches(0.5))
    runs_para(tf, [("paper-writer", ORANGE, True, 14), (" 出初稿  →  ", GREY, False, 13.5),
                   ("paper-pipeline", ORANGE, True, 14),
                   (" 按固定顺序调度 5 个技能，顺序本身就是方法论。", GREY, False, 13.5)],
              first=True, space_after=0)
    stages = [
        ("1", "paper-polish", "机械/一致性/引用纠错", TEAL),
        ("2", "paper-self-revise", "在干净稿上做深度内容自审", GOLD),
        ("3", "paper-style", "按目标期刊重构标题与结构", ORANGE),
        ("4", "paper-polish", "二次润色：清理 2–3 步引入的接缝", TEAL),
        ("5", "reference-verify", "最后做引用核查（生成 xlsx 报告）", NAVY),
    ]
    x = Inches(0.75); y = Inches(2.6); bw = Inches(2.18); bh = Inches(2.4); gap = Inches(0.18)
    for i,(n, name, desc, ac) in enumerate(stages):
        rect(s, x, y, bw, bh, CARD, line_color=LINE)
        rect(s, x, y, bw, Inches(0.55), ac)
        tb, tf = textbox(s, x, y + Inches(0.03), bw, Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
        para(tf, f"Stage {n}", size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
             first=True, space_after=0)
        tb, tf = textbox(s, x + Inches(0.12), y + Inches(0.68), bw - Inches(0.24), Inches(1.6))
        para(tf, name, size=13.5, color=NAVY, bold=True, first=True, space_after=6, line_spacing=1.05)
        para(tf, desc, size=11.5, color=GREY, space_after=0, line_spacing=1.2)
        if i < len(stages)-1:
            arrow(s, x + bw - Inches(0.04), y + Inches(0.9), Inches(0.2), Inches(0.55), ORANGE)
        x += bw + gap
    # 底注: 为什么是这个顺序
    y2 = Inches(5.35)
    rect(s, Inches(0.75), y2, Inches(11.85), Inches(1.2), LGREY, line_color=LINE)
    tb, tf = textbox(s, Inches(1.05), y2 + Inches(0.13), Inches(11.3), Inches(0.95),
                     anchor=MSO_ANCHOR.MIDDLE)
    runs_para(tf, [("为什么顺序不能乱：", NAVY, True, 13.5),
                   ("先润色让自审聚焦实质；内容稳定后再改结构避免反复；"
                    "因为 2–3 步动了文字，所以要二次润色；引用核查放最后，否则会被后续编辑作废。",
                    GREY, False, 13)], first=True, space_after=0, line_spacing=1.25)


# ============================================================================
# SLIDE 27 — 评审与引用核查
# ============================================================================
def slide_review():
    s = add_slide()
    content_header(s, "7 · 评审", "在投出去之前，先让 AI 当一次审稿人", 27)
    cards3 = [
        ("referee-report", ORANGE, "模拟审稿",
         ["以目标期刊审稿人视角写报告", "指出识别威胁、稳健性缺口", "区分『大修 / 小修』级问题"]),
        ("paper-referee-revise", TEAL, "回应修改",
         ["逐条回应审稿意见", "生成 response letter 草稿", "把修改落到正文对应位置"]),
        ("reference-verify", NAVY, "引用核查",
         ["核对每条引用真实存在", "检查作者/年份/期刊一致", "输出 xlsx 差错报告"]),
    ]
    x = Inches(0.75); y = Inches(2.0); cw = Inches(3.86); gap = Inches(0.13)
    for name, ac, role, lines in cards3:
        card(s, x, y, cw, Inches(2.95), name, lines, accent=ac, tag=role,
             title_size=15, body_size=12.3)
        x += cw + gap
    # 底: master-thesis-review & novelty 复用
    y2 = Inches(5.25)
    rect(s, Inches(0.75), y2, Inches(11.85), Inches(1.3), NAVY)
    rect(s, Inches(0.75), y2, Inches(0.12), Inches(1.3), ORANGE)
    tb, tf = textbox(s, Inches(1.1), y2 + Inches(0.15), Inches(11.2), Inches(1.05),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "评审不是终点，是又一轮迭代的起点", size=14.5, color=ORANGE, bold=True,
         first=True, space_after=5)
    para(tf, "审稿意见若涉及新的实证（如补一个稳健性），回到第④阶段重跑 did-analysis / 相应估计技能 —— "
             "工作流天然支持『回环』。", size=12.5, color=RGBColor(0xC6,0xD2,0xE2),
         space_after=0, line_spacing=1.22)


# ============================================================================
# SLIDE 28 — 投稿与期刊匹配
# ============================================================================
def slide_submission():
    s = add_slide()
    content_header(s, "7 · 投稿", "paper-submission：把『一篇论文』打包成『一次投稿』", 28)
    left = [
        ("期刊匹配", TEAL, ["按主题/方法/样本匹配候选期刊", "对照 journal-digest 画像", "排序：契合度 × 命中率"]),
        ("格式适配", GOLD, ["按目标期刊模板排版", "字数、章节、参考文献样式", "图表标题与编号规范"]),
        ("投稿材料", ORANGE, ["cover letter 草稿", "highlights / 摘要要点", "推荐审稿人 + 利益冲突声明"]),
    ]
    x = Inches(0.75); y = Inches(2.0); cw = Inches(3.86); gap = Inches(0.13)
    for name, ac, lines in left:
        card(s, x, y, cw, Inches(2.6), name, lines, accent=ac, title_size=15.5, body_size=12.3)
        x += cw + gap
    # 底: 期刊池说明
    y2 = Inches(4.85)
    rect(s, Inches(0.75), y2, Inches(11.85), Inches(1.7), CARD, line_color=LINE)
    rect(s, Inches(0.75), y2, Inches(0.12), Inches(1.7), NAVY)
    tb, tf = textbox(s, Inches(1.1), y2 + Inches(0.18), Inches(11.2), Inches(1.4))
    para(tf, "工具箱内置的期刊知识", size=14.5, color=NAVY, bold=True, first=True, space_after=6)
    runs_para(tf, [("• ", ORANGE, True, 12.5),
                   ("结构化期刊清单（journals.json）涵盖 60+ 英文顶刊（经济/金融/会计）；", INK, False, 12.5)],
              space_after=4)
    runs_para(tf, [("• ", ORANGE, True, 12.5),
                   ("China-CF-study / Foreign-CF-study 提供中国情境与海外公司金融的目标期刊池与写作范式；", INK, False, 12.5)],
              space_after=4)
    runs_para(tf, [("• ", ORANGE, True, 12.5),
                   ("paper-style 据此把同一篇稿子，快速改写成不同期刊的『口味』。", INK, False, 12.5)],
              space_after=0)


# ============================================================================
# SLIDE 29 — 全流程复盘
# ============================================================================
def slide_recap():
    s = add_slide()
    content_header(s, "复盘", "一页看懂：技能如何串成一条研究流水线", 29)
    rows = [
        ("①选题", "idea-finder · novelty-check · significance-search", "研究问题 + 贡献", TEAL),
        ("②设计", "econfin-proposal · journal-digest", "识别策略 + 期刊画像", TEAL),
        ("③数据", "data-fetcher · data-cleaning", "可复现分析数据集", GOLD),
        ("④估计", "did-analysis · panel-data · iv · rdd · scm · stata", "ATT + 稳健性", ORANGE),
        ("⑤呈现", "table · figure", "出版级表 & 图", GOLD),
        ("⑥写作", "paper-writer · paper-pipeline（polish→revise→style）", "结构完整定稿", TEAL),
        ("⑦评审", "referee-report · paper-referee-revise · reference-verify", "意见 + 回应", TEAL),
        ("⑧投稿", "paper-submission", "投稿包 + cover letter", NAVY),
    ]
    x0 = Inches(0.75); y0 = Inches(1.95); rh = Inches(0.56)
    c1 = Inches(1.5); c2 = Inches(7.4); c3 = Inches(2.95)
    # 表头
    rect(s, x0, y0, c1+c2+c3, Inches(0.45), NAVY)
    for cx, w, h in [(x0, c1, "阶段"), (x0+c1, c2, "调用的技能"), (x0+c1+c2, c3, "产出物")]:
        tb, tf = textbox(s, cx+Inches(0.12), y0, w-Inches(0.2), Inches(0.45), anchor=MSO_ANCHOR.MIDDLE)
        para(tf, h, size=12.5, color=WHITE, bold=True, first=True, space_after=0)
    for i,(st, sk, out, ac) in enumerate(rows):
        y = y0 + Inches(0.45) + i*rh
        fill = WHITE if i%2==0 else LGREY
        rect(s, x0, y, c1+c2+c3, rh, fill, line_color=LINE)
        rect(s, x0, y, Inches(0.07), rh, ac)
        tb, tf = textbox(s, x0+Inches(0.14), y, c1-Inches(0.2), rh, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, st, size=12.5, color=ac, bold=True, first=True, space_after=0)
        tb, tf = textbox(s, x0+c1+Inches(0.1), y, c2-Inches(0.2), rh, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, sk, size=11.5, color=INK, first=True, space_after=0, line_spacing=1.0)
        tb, tf = textbox(s, x0+c1+c2+Inches(0.1), y, c3-Inches(0.2), rh, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, out, size=11.5, color=GREY, bold=True, first=True, space_after=0, line_spacing=1.0)


# ============================================================================
# SLIDE 30 — 结束页
# ============================================================================
def slide_end():
    s = add_slide()
    bg(s, NAVY)
    rect(s, Inches(9.6), Inches(0), Inches(3.733), EMU_H, NAVY_DK)
    rect(s, Inches(9.6), Inches(0), Inches(0.06), EMU_H, ORANGE)
    tb, tf = textbox(s, Inches(0.9), Inches(1.6), Inches(8.4), Inches(2.2))
    para(tf, "把流程交给工作流，", size=40, color=WHITE, bold=True, first=True, space_after=2)
    para(tf, "把判断留给自己。", size=40, color=ORANGE, bold=True, space_after=16)
    para(tf, "AI 负责跑通每一步的脏活，研究者负责识别、解释与取舍。",
         size=15, color=RGBColor(0xC6,0xD2,0xE2), space_after=0, line_spacing=1.3)
    rect(s, Inches(0.95), Inches(4.2), Inches(2.6), Pt(3), ORANGE)
    # 资源区
    tb, tf = textbox(s, Inches(0.9), Inches(4.5), Inches(8.4), Inches(2.2))
    para(tf, "本次交付物（skills/Paper-WorkFlow/）", size=15, color=ORANGE, bold=True,
         first=True, space_after=8)
    for ln in ["• 社科实证论文工作流.pptx   —   本 30 页演示文稿",
               "• did_demo.ipynb   —   可一键运行的 DiD 演示 Notebook",
               "• build_pptx.py   —   PPTX 生成脚本（可改主题/内容重生成）",
               "• README.md   —   使用说明与技能映射"]:
        para(tf, ln, size=13.5, color=WHITE, space_after=6)
    # 右侧
    tb, tf = textbox(s, Inches(10.0), Inches(2.7), Inches(3.0), Inches(2.5))
    para(tf, "演示方法", size=13, color=RGBColor(0x9F,0xB0,0xC8), bold=True, first=True, space_after=4)
    para(tf, "双重差分 DiD", size=22, color=WHITE, bold=True, space_after=14)
    para(tf, "技能来源", size=13, color=RGBColor(0x9F,0xB0,0xC8), bold=True, space_after=4)
    para(tf, "econfin-\nworkflow-toolkit", size=16, color=WHITE, bold=True, space_after=14, line_spacing=1.1)
    para(tf, "47 skills · 8 stages", size=12, color=ORANGE, bold=True, space_after=0)


# ============================================================================
# 组装
# ============================================================================
def build():
    slide_cover()                 # 1
    slide_agenda()                # 2
    slide_section("01", "为什么需要『工作流』", "WHY WORKFLOW",
                  "实证研究真正的难点在『串联』——\n把选题、数据、识别、写作、投稿连成一条可复现、可交接的主线。", 3)  # 3
    slide_pain()                  # 4
    slide_pipeline_overview()     # 5
    slide_toolkit_map()           # 6
    slide_section("02", "选题与研究设计", "IDEATION & DESIGN",
                  "把模糊的兴趣，磨成一个有新意、可识别、能投稿的研究问题，\n并翻译成可执行的研究计划书。", 7)  # 7
    slide_ideation()              # 8
    slide_proposal()              # 9
    slide_section("03", "数据获取与清洗", "DATA",
                  "定位来源 → 抓取 → 对齐键 → 合并 → 清洗 → 留痕。\n这一阶段的产出，是一张可复现的『分析主表』。", 10)  # 10
    slide_data_fetch()            # 11
    slide_data_clean()            # 12
    slide_section("04", "计量识别与估计", "ESTIMATION",
                  "方法服从于识别：先问『反事实从哪来』，\n再从八种识别策略里选对应的估计量。", 13)  # 13
    slide_method_map()            # 14
    slide_decision_tree()         # 15
    slide_section("05", "聚焦 DiD（本次演示）", "FOCUS · DiD",
                  "用『差中之差』剥离共同趋势 —— 从核心逻辑、平行趋势、\n交错偏误到稳健性，配一个可运行的 Notebook。", 16)  # 16
    slide_did_core()              # 17
    slide_event_study()           # 18
    slide_staggered()             # 19
    slide_robustness()            # 20
    slide_notebook()              # 21
    slide_section("06", "结果呈现", "TABLES & FIGURES",
                  "把估计结果排成审稿人一眼能读的表，\n把识别直觉画成一张胜过半页文字的图。", 22)  # 22
    slide_table()                 # 23
    slide_figure()                # 24
    slide_section("07", "写作 · 评审 · 投稿", "WRITE · REVIEW · SUBMIT",
                  "初稿之后的一条龙打磨、内部模拟审稿与引用核查，\n最后打包成一次面向目标期刊的投稿。", 25)  # 25
    slide_writing()               # 26
    slide_review()                # 27
    slide_submission()            # 28
    slide_recap()                 # 29
    slide_end()                   # 30

    out = "社科实证论文工作流.pptx"
    prs.save(out)
    print(f"已生成: {out}  共 {len(prs.slides.__iter__.__self__._sldIdLst)} 页")


if __name__ == "__main__":
    build()
