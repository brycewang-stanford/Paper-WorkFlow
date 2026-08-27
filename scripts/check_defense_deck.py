#!/usr/bin/env python3
"""Defence-deck generator gate for Paper-WorkFlow.

Why this exists
---------------
`references/stage-playbook.md` Stage 9 calls the defence deck a **hard
deliverable** for the Chinese thesis path — "不能因脚本报错就跳过答辩 PPT". That
made `defense_pptx.py` the largest load-bearing artifact in the repository with
no executable gate behind it, in a package whose entire thesis is that rigor is
executable rather than advisory.

It was also wrong in a way only execution finds. The thesis template rationed a
fixed 22-slide budget in first-come order, so a run with six findings consumed
the budget before reaching the closing sections and the deck shipped with **no
研究结论与主要贡献, no 机制分析, and no 研究局限与展望** — the three slides a
defence committee is guaranteed to ask about. Nothing reported it; the generator
printed a success line and exited 0.

The invariant this gate enforces:

    The mandatory closing sections are never rationed away. A variable-length
    findings list is what gets truncated, and truncation is announced.

Plus the basics that keep the deliverable a deliverable: the generator runs, the
deck opens, the slide count respects the template envelope, `--duration`
actually changes the budget, and workspace content reaches the slides rather
than leaving placeholders behind.

Usage:
    python3 check_defense_deck.py            # build + verify decks in a temp workspace
    python3 check_defense_deck.py --json
    python3 check_defense_deck.py --selftest # same, plus the pure-logic invariants

Requires python-dev extras (`pip install -r requirements-dev.txt`); a missing
python-pptx is a hard failure, exactly as the missing notebook stack is for
check_demo_execution.py — an ungenerable hard deliverable is not a pass.
"""

from __future__ import annotations

import argparse
import contextlib
import importlib.util
import io
import json
import subprocess
import sys
import tempfile
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
GENERATOR = ROOT / "defense_pptx.py"

# Sections a Chinese thesis defence deck may not ship without. These are the
# ones the fixed-budget bug used to drop.
THESIS_REQUIRED_SLIDES = [
    "选题背景",
    "文献综述与本文定位",
    "理论分析与研究假设",
    "数据与样本",
    "识别策略",
    "主要发现",
    "稳健性检验",
    "机制分析与异质性",
    "研究结论与主要贡献",
    "研究局限与未来展望",
]
JOURNAL_TALK_REQUIRED_SLIDES = [
    "研究问题",
    "文献定位",
    "数据描述",
    "Main Finding",
    "Robustness",
    "Contribution",
]
# (type, slides at the documented default duration, envelope)
TEMPLATE_ENVELOPE = {"thesis": (22, 18, 28), "journal-talk": (18, 15, 22)}


def _fail(msg: str) -> None:
    print(f"FAIL: {msg}", file=sys.stderr)
    raise SystemExit(1)


def _load_generator():
    name = "_defense_pptx"
    if name in sys.modules:
        return sys.modules[name]
    spec = importlib.util.spec_from_file_location(name, GENERATOR)
    module = importlib.util.module_from_spec(spec)
    # @dataclass resolves annotations through sys.modules, so register before exec.
    sys.modules[name] = module
    try:
        spec.loader.exec_module(module)
    except BaseException:
        sys.modules.pop(name, None)
        raise
    return module


@contextlib.contextmanager
def _quiet():
    """The generator narrates to stdout/stderr; the gate reports its own verdict."""
    buf = io.StringIO()
    with contextlib.redirect_stdout(buf), contextlib.redirect_stderr(buf):
        yield buf


def _build(module, cfg, out: Path) -> str:
    """Build a deck, returning whatever the generator narrated."""
    with _quiet() as buf:
        if cfg.type == "thesis":
            module.build_thesis_pptx(cfg, out)
        else:
            module.build_journal_talk_pptx(cfg, out)
    return buf.getvalue()


def _slide_texts(path: Path) -> list[str]:
    from pptx import Presentation
    prs = Presentation(str(path))
    out = []
    for slide in prs.slides:
        parts = [sh.text_frame.text for sh in slide.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        out.append("\n".join(parts))
    return out


def build_workspace(root: Path, n_findings: int) -> Path:
    """A workspace filled enough that the generator has real content to place."""
    ws = root / "paper_workspace" / "deck_project"
    subprocess.run(["bash", str(ROOT / "assets" / "init_workspace.sh"), str(ws)],
                   check=True, capture_output=True)
    (ws / "00_meta" / "intake.md").write_text(
        "# Intake\n- 论文题目：绿色信贷政策与企业创新\n- 学生：王磊\n"
        "- 导师：张伟\n- 学校：示例大学\n- 学院：经济学院\n- 学位：硕士\n",
        encoding="utf-8")
    (ws / "01_proposal").mkdir(exist_ok=True)
    (ws / "01_proposal" / "proposal.md").write_text(
        "# 绿色信贷政策与企业创新\n\n"
        "## 一、选题背景与意义\n- 绿色信贷自 2012 年起在全国推开\n- 融资约束是创新的主要瓶颈\n\n"
        "## 二、研究意义\n- 为绿色金融的实体效应提供微观证据\n\n"
        "## 三、文献综述\n- 环境规制与创新的波特假说文献\n- 信贷配置与企业投资文献\n\n"
        "## 四、理论分析\n- 融资成本上升挤出污染型投资\n",
        encoding="utf-8")
    (ws / "03_analysis" / "results").mkdir(parents=True, exist_ok=True)
    (ws / "03_analysis" / "results" / "main_results.json").write_text(
        json.dumps({"estimates": [
            {"variable": f"treat_x{i}", "coefficient": 0.100 + i / 100,
             "std_error": 0.021, "p_value": 0.01, "title": f"主发现 {i}"}
            for i in range(1, n_findings + 1)]}, ensure_ascii=False),
        encoding="utf-8")
    return ws


def check(verbose: bool = True) -> list[dict]:
    try:
        import pptx  # noqa: F401
    except ImportError:
        _fail("defence-deck gate needs python-pptx; install with: "
              "pip install -r requirements-dev.txt")

    module = _load_generator()
    results: list[dict] = []

    with tempfile.TemporaryDirectory(prefix="defense-deck-") as tmp:
        root = Path(tmp)

        # --- 1. many findings must NOT crowd out the mandatory closing slides
        ws = build_workspace(root, n_findings=6)
        for deck_type, required in (("thesis", THESIS_REQUIRED_SLIDES),
                                    ("journal-talk", JOURNAL_TALK_REQUIRED_SLIDES)):
            out = root / f"{deck_type}.pptx"
            cfg = module.DefenseConfig.from_workspace(ws)
            cfg.type = deck_type
            cfg.duration_min = 15
            _build(module, cfg, out)
            if not out.exists():
                _fail(f"{deck_type}: generator produced no file")
            texts = _slide_texts(out)
            base, lo, hi = TEMPLATE_ENVELOPE[deck_type]
            if not lo <= len(texts) <= hi:
                _fail(f"{deck_type}: {len(texts)} slides is outside the template envelope {lo}-{hi}")
            blob = "\n".join(texts)
            missing = [s for s in required if s not in blob]
            if missing:
                _fail(f"{deck_type}: six findings crowded out mandatory section(s): {missing}. "
                      "The findings list is what gets truncated, never the conclusion.")
            # the closing sections must come after the findings, not merely exist
            if deck_type == "thesis":
                last_finding = max(i for i, t in enumerate(texts) if "主要发现" in t)
                conclusion = next(i for i, t in enumerate(texts) if "研究结论与主要贡献" in t)
                if conclusion < last_finding:
                    _fail("thesis: 研究结论与主要贡献 appears before the findings")
            results.append({"case": f"{deck_type}_six_findings", "slides": len(texts), "ok": True})

        # --- 2. workspace content actually reaches the deck
        cfg = module.DefenseConfig.from_workspace(ws)
        cfg.type = "thesis"
        out = root / "content.pptx"
        _build(module, cfg, out)
        blob = "\n".join(_slide_texts(out))
        for needle in ("绿色信贷政策与企业创新", "王磊", "绿色信贷自 2012 年起在全国推开", "0.1"):
            if needle not in blob:
                _fail(f"workspace content did not reach the deck: {needle!r} missing")
        results.append({"case": "content_extraction", "ok": True})

        # --- 3. --duration actually changes the budget (the hint promises it does)
        cfg = module.DefenseConfig.from_workspace(ws)
        cfg.type = "thesis"
        cfg.duration_min = 25
        long_out = root / "long.pptx"
        _build(module, cfg, long_out)
        long_texts = _slide_texts(long_out)
        short_texts = _slide_texts(root / "thesis.pptx")
        if len(long_texts) <= len(short_texts):
            _fail(f"--duration had no effect: 25min={len(long_texts)} slides, "
                  f"15min={len(short_texts)}")
        if sum("主要发现" in t for t in long_texts) < 6:
            _fail("a longer defence must fit all six findings")
        for needle in ("研究结论与主要贡献", "研究局限与未来展望"):
            if needle not in "\n".join(long_texts):
                _fail(f"long deck lost mandatory section: {needle}")
        results.append({"case": "duration_scales_budget",
                        "slides_15min": len(short_texts), "slides_25min": len(long_texts),
                        "ok": True})

        # --- 4. an empty workspace still produces a complete, usable deck
        empty = build_workspace(root / "empty", n_findings=0)
        (empty / "03_analysis" / "results" / "main_results.json").unlink()
        cfg = module.DefenseConfig.from_workspace(empty)
        cfg.type = "thesis"
        out = root / "empty.pptx"
        _build(module, cfg, out)
        texts = _slide_texts(out)
        blob = "\n".join(texts)
        missing = [s for s in THESIS_REQUIRED_SLIDES if s not in blob]
        if missing:
            _fail(f"empty workspace lost mandatory section(s): {missing}")
        results.append({"case": "empty_workspace", "slides": len(texts), "ok": True})

        # --- 5. the CLI path (what the playbook actually tells people to run)
        cli_out = root / "cli.pptx"
        proc = subprocess.run(
            [sys.executable, str(GENERATOR), "--workspace", str(ws),
             "--type", "thesis", "--duration", "15", "--output", str(cli_out)],
            capture_output=True, text=True)
        if proc.returncode != 0:
            _fail(f"documented CLI invocation failed: {proc.stderr.strip()}")
        if not cli_out.exists():
            _fail("documented CLI invocation exited 0 without producing a deck")
        if "只放前" not in proc.stderr:
            _fail("findings were truncated without saying so — silent truncation is the "
                  "failure mode this gate exists to prevent")
        results.append({"case": "cli_invocation", "ok": True})

    if verbose:
        print("Paper-WorkFlow defence-deck gate")
        for r in results:
            detail = " ".join(f"{k}={v}" for k, v in r.items() if k not in {"case", "ok"})
            print(f"  [OK] {r['case']}{(' ' + detail) if detail else ''}")
        print("  DEFENCE DECK OK")
    return results


def _selftest() -> int:
    module = _load_generator()

    # The budget helper is the whole fix; test it without building decks.
    budget = module._findings_budget
    _b = budget
    def budget(**kw):
        with _quiet():
            return _b(**kw)
    assert budget(total=22, fixed_slides=19, wanted=6, hard_cap=6) == 3, "must reserve the closing sections"
    assert budget(total=28, fixed_slides=19, wanted=6, hard_cap=6) == 6, "a longer deck fits them all"
    assert budget(total=22, fixed_slides=19, wanted=2, hard_cap=6) == 2, "never pad beyond what exists"
    assert budget(total=19, fixed_slides=19, wanted=4, hard_cap=6) == 1, "always at least one findings slide"
    assert budget(total=28, fixed_slides=19, wanted=99, hard_cap=6) == 6, "hard cap still applies"
    # ...and truncation is never silent (use the raw helper, not the quiet wrapper)
    with _quiet() as buf:
        _b(total=22, fixed_slides=19, wanted=6, hard_cap=6)
    assert "只放前 3 条" in buf.getvalue(), buf.getvalue()
    with _quiet() as buf:
        _b(total=28, fixed_slides=19, wanted=6, hard_cap=6)
    assert buf.getvalue() == "", "a deck that fits everything must not warn"

    # duration -> slide budget, anchored on the documented defaults so the
    # historical behaviour at 15 minutes is unchanged.
    cfg = module.DefenseConfig(title="t")
    for deck_type, expected_default in (("thesis", 22), ("journal-talk", 18)):
        cfg.type = deck_type
        cfg.duration_min = 15
        assert cfg.total_slides == expected_default, (deck_type, cfg.total_slides)
        cfg.duration_min = 30
        long_total = cfg.total_slides
        cfg.duration_min = 5
        short_total = cfg.total_slides
        assert short_total < expected_default < long_total, (deck_type, short_total, long_total)
        lo, hi = module_envelope(deck_type)
        assert lo <= short_total and long_total <= hi, (deck_type, short_total, long_total)

    check(verbose=False)
    print("selftest OK: defence-deck invariants hold")
    return 0


def module_envelope(deck_type: str) -> tuple[int, int]:
    _, lo, hi = TEMPLATE_ENVELOPE[deck_type]
    return lo, hi


def main(argv: list[str] | None = None) -> int:
    p = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    p.add_argument("--json", action="store_true", help="machine-readable output")
    p.add_argument("--selftest", action="store_true", help="pure-logic invariants + the full build")
    args = p.parse_args(argv)
    if args.selftest:
        return _selftest()
    results = check(verbose=not args.json)
    if args.json:
        print(json.dumps({"ok": True, "cases": results}, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
